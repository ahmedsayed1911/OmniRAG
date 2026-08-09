"""Shared pytest fixtures.

Two rules the whole suite follows:

* **no paid API keys** — every external provider is faked or mocked;
* **no network** — nothing here opens a socket.
"""

from __future__ import annotations

import os
import sys
from pathlib import Path
from typing import Dict, List, Optional, Sequence

import pytest

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

# Env vars that must not leak from the developer's shell into the tests.
_PROVIDER_ENV = [
    "LLM_PROVIDER", "LLM_API_KEY", "LLM_MODEL", "LLM_BASE_URL",
    "LLM_MAX_OUTPUT_TOKENS", "LLM_EXHAUSTIVE_MAX_OUTPUT_TOKENS",
    "PRIMARY_LLM_PROVIDER", "FALLBACK_LLM_PROVIDER", "ENABLE_PROVIDER_FALLBACK",
    "GEMINI_API_KEY", "GEMINI_MODEL", "GOOGLE_API_KEY",
    "OPENROUTER_API_KEY", "OPENROUTER_MODEL", "OPENROUTER_MODEL_SUPPORTS_IMAGES",
    "OPENAI_API_KEY", "ANTHROPIC_API_KEY", "COHERE_API_KEY", "JINA_API_KEY",
    "EMBEDDING_PROVIDER", "EMBEDDING_API_KEY", "EMBEDDING_MODEL",
    "QDRANT_URL", "QDRANT_API_KEY", "QDRANT_COLLECTION",
    "RERANK_PROVIDER", "RERANK_API_KEY", "RERANK_MODEL",
    "OCR_PROVIDER", "VISION_ENABLED", "QUERY_REWRITE",
    "TOP_K", "RERANK_TOP_K", "EXHAUSTIVE_SCAN_MAX_CHUNKS", "EXHAUSTIVE_FINAL_K",
    "CHUNK_SIZE", "CHUNK_OVERLAP", "MAX_UPLOAD_MB",
    "VECTOR_STORE_PROVIDER", "DEBUG_PANELS", "LOG_LEVEL",
]


@pytest.fixture(autouse=True)
def clean_environment(monkeypatch):
    """Start every test from a known, key-free environment."""
    for name in _PROVIDER_ENV:
        monkeypatch.delenv(name, raising=False)

    from omnirag.config.settings import reset_settings_cache

    reset_settings_cache()
    yield
    reset_settings_cache()


@pytest.fixture(autouse=True)
def reset_singletons():
    """Isolate module-level caches between tests."""
    from omnirag.ingestion.router import reset_router
    from omnirag.providers.embeddings.factory import reset_embedding_cache
    from omnirag.providers.llm.factory import reset_llm_cache
    from omnirag.providers.ocr.factory import reset_ocr_cache
    from omnirag.providers.rerank.factory import reset_rerank_cache
    from omnirag.rag.hybrid import get_bm25_cache
    from omnirag.rag.vector_store import set_vector_store
    from omnirag.services.engine import reset_engine
    from omnirag.storage.sessions import reset_registry

    def _reset():
        reset_llm_cache()
        reset_embedding_cache()
        reset_rerank_cache()
        reset_ocr_cache()
        reset_registry()
        reset_router()
        reset_engine()
        set_vector_store(None)
        get_bm25_cache().clear()

    _reset()
    yield
    _reset()


# --------------------------------------------------------------------------- #
@pytest.fixture
def settings():
    """Default settings with the offline providers selected."""
    from omnirag.config.settings import build_settings

    os.environ["PRIMARY_LLM_PROVIDER"] = "mock"
    os.environ["FALLBACK_LLM_PROVIDER"] = "none"
    os.environ["EMBEDDING_PROVIDER"] = "hash"
    os.environ["VISION_ENABLED"] = "false"
    os.environ["QUERY_REWRITE"] = "false"
    try:
        yield build_settings()
    finally:
        for name in (
            "PRIMARY_LLM_PROVIDER", "FALLBACK_LLM_PROVIDER",
            "EMBEDDING_PROVIDER", "VISION_ENABLED", "QUERY_REWRITE",
        ):
            os.environ.pop(name, None)


@pytest.fixture
def session_id() -> str:
    from omnirag.storage.sessions import new_session_id

    return new_session_id()


@pytest.fixture
def file_store():
    from omnirag.storage.files import MemoryFileStore

    return MemoryFileStore()


@pytest.fixture
def vector_store():
    from omnirag.rag.vector_store import InMemoryVectorStore

    return InMemoryVectorStore()


@pytest.fixture
def engine(settings, file_store, vector_store):
    """Engine wired entirely with offline components."""
    from omnirag.services.engine import OmniRAGEngine

    built = OmniRAGEngine(settings)
    built._file_store = file_store
    built._vector_store = vector_store
    yield built


@pytest.fixture
def ingestion_service(engine):
    from omnirag.services.ingestion_service import IngestionService

    return IngestionService(engine)


# --------------------------------------------------------------------------- #
# Sample documents
# --------------------------------------------------------------------------- #
SAMPLE_MARKDOWN = """# Annual Report 2024

## Revenue Overview

Total revenue reached 8,400,000 USD in Q4 2024, up from 6,200,000 USD in Q3 2024.
Growth was driven primarily by the enterprise segment in EMEA.

## Regional Breakdown

| Region | Q3 2024 | Q4 2024 |
| --- | --- | --- |
| EMEA | 2100 | 3400 |
| APAC | 1800 | 2600 |
| Americas | 2300 | 2400 |

## Risks

Currency volatility remains the principal risk to the 2025 outlook.

## الملخص التنفيذي

بلغت الإيرادات الإجمالية 8.4 مليون دولار في الربع الرابع من عام 2024.
ويظل تقلب أسعار الصرف هو الخطر الرئيسي على توقعات عام 2025.
"""

SAMPLE_TEXT = """Project Atlas Status Report

The migration completed on 12 March 2024 with zero downtime.
Total cost was 145,000 EUR against a budget of 160,000 EUR.

Outstanding items: documentation refresh and load testing.
"""


@pytest.fixture
def sample_markdown() -> bytes:
    return SAMPLE_MARKDOWN.encode("utf-8")


@pytest.fixture
def sample_text() -> bytes:
    return SAMPLE_TEXT.encode("utf-8")


@pytest.fixture
def sample_png() -> bytes:
    """A small non-uniform PNG (passes the decorative-image filter)."""
    from PIL import Image, ImageDraw

    import io

    image = Image.new("RGB", (420, 320), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle([20, 20, 200, 160], fill="#3355aa")
    draw.rectangle([220, 60, 400, 300], fill="#aa5533")
    draw.line([0, 0, 420, 320], fill="black", width=3)
    draw.text((30, 200), "Revenue 8400", fill="black")
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


@pytest.fixture
def sample_pdf() -> bytes:
    """A two-page digital PDF built at test time (no binary fixtures in git)."""
    pymupdf = pytest.importorskip("pymupdf")

    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 100), "Quarterly Report 2024", fontsize=20)
    page.insert_text((72, 140), "Total revenue reached 8,400,000 USD in Q4 2024.", fontsize=11)
    page.insert_text((72, 160), "The enterprise segment drove the increase.", fontsize=11)

    page2 = doc.new_page()
    page2.insert_text((72, 100), "Risk Factors", fontsize=20)
    page2.insert_text((72, 140), "Currency volatility remains the principal risk.", fontsize=11)

    data = doc.tobytes()
    doc.close()
    return data


@pytest.fixture
def sample_docx() -> bytes:
    """A DOCX with a heading, paragraphs and a table."""
    docx = pytest.importorskip("docx")

    import io

    document = docx.Document()
    document.add_heading("Operations Review", level=1)
    document.add_paragraph("Throughput increased by 18% during the second half of 2024.")
    document.add_heading("Metrics", level=2)
    table = document.add_table(rows=3, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Throughput"
    table.cell(1, 1).text = "1840"
    table.cell(2, 0).text = "Defects"
    table.cell(2, 1).text = "12"

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


@pytest.fixture
def sample_pptx() -> bytes:
    """A PPTX with a title slide, bullets, a table and speaker notes."""
    pptx = pytest.importorskip("pptx")

    import io

    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Market Expansion"
    slide.placeholders[1].text = "EMEA grew 62%\nAPAC grew 44%"
    slide.notes_slide.notes_text_frame.text = "Mention the currency headwind."

    slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide2.shapes.title.text = "Numbers"
    table = slide2.shapes.add_table(
        3, 2, Inches(1), Inches(2), Inches(6), Inches(2)
    ).table
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Growth"
    table.cell(1, 0).text = "EMEA"
    table.cell(1, 1).text = "62%"
    table.cell(2, 0).text = "APAC"
    table.cell(2, 1).text = "44%"

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


# --------------------------------------------------------------------------- #
# Fakes
# --------------------------------------------------------------------------- #
class FakeEmbeddings:
    """Deterministic embeddings with no API calls.

    Token-overlap based, so semantically similar sentences really do score
    higher — enough to make retrieval assertions meaningful.
    """

    name = "fake"
    dimensions = 64
    batch_size = 32
    supports_task_type = False
    max_chars = 8000

    def __init__(self, dimensions: int = 64):
        self.dimensions = dimensions
        self.calls = 0

    def embed_batch(self, texts, *, is_query: bool = False):
        return [self._vector(t) for t in texts]

    def embed_documents(self, texts: Sequence[str]) -> List[List[float]]:
        self.calls += 1
        return [self._vector(t) for t in texts]

    def embed_query(self, text: str) -> List[float]:
        self.calls += 1
        return self._vector(text)

    def _vector(self, text: str) -> List[float]:
        import math

        from omnirag.utils.text import tokenize

        vector = [0.0] * self.dimensions
        for token in tokenize(text):
            vector[hash(token) % self.dimensions] += 1.0
        norm = math.sqrt(sum(v * v for v in vector)) or 1.0
        return [v / norm for v in vector]


@pytest.fixture
def fake_embeddings():
    return FakeEmbeddings()


class RecordingLLM:
    """Scripted LLM: returns queued responses, records the calls it received."""

    name = "recording"
    supports_vision = True

    def __init__(self, responses: Optional[List[str]] = None, model: str = "recording-1"):
        self.model = model
        self.temperature = 0.0
        self.max_output_tokens = 1000
        self.responses = list(responses or [])
        self.calls: List[Dict] = []

    def supports_images(self, model: Optional[str] = None) -> bool:
        return True

    def complete(self, messages, *, system=None, temperature=None,
                 max_output_tokens=None, model=None, json_mode=False):
        from omnirag.providers.llm.base import LLMResponse

        self.calls.append(
            {
                "system": system,
                "text": "\n".join(m.text for m in messages),
                "images": sum(len(m.images) for m in messages),
                "json_mode": json_mode,
                "roles": [str(m.role) for m in messages],
            }
        )
        text = self.responses.pop(0) if self.responses else "No answer available."
        return LLMResponse(text=text, model=self.model, provider=self.name)


@pytest.fixture
def recording_llm():
    return RecordingLLM()
