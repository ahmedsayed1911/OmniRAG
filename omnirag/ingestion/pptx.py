"""PPTX processing with python-pptx.

A slide maps directly onto a page, so citations read "Slide 7" — which is
exactly what the user sees in PowerPoint.

Extracted per slide: the title, body text (following the shape order), tables,
speaker notes, pictures, and grouped shapes (recursively). Charts embedded as
native PowerPoint objects expose their plot data through the XML, so those
numbers are read structurally rather than guessed from a picture; charts that
are only images go through the visual pipeline. Diagrams built from grouped
autoshapes cannot be exported as an image by python-pptx, so their text and
relationships are captured from the shapes themselves.
"""

from __future__ import annotations

import io
from typing import Any, Iterable, List, Optional

from omnirag.core.enums import BlockType, FileType, PipelineStage, SourceKind
from omnirag.core.exceptions import CorruptedDocumentError, EmptyDocumentError
from omnirag.core.models import ContentBlock, Document, Page
from omnirag.ingestion.base import BaseDocumentProcessor, ProcessingContext
from omnirag.intelligence.tables import build_table, table_to_text
from omnirag.utils.images import is_probably_decorative
from omnirag.utils.logging import get_logger
from omnirag.utils.text import clean_text

logger = get_logger(__name__)

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    PYTHON_PPTX_AVAILABLE = True
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore[assignment]
    MSO_SHAPE_TYPE = None  # type: ignore[assignment]
    PYTHON_PPTX_AVAILABLE = False

MIN_IMAGE_BYTES = 2048
MAX_IMAGES_PER_SLIDE = 4


class PowerPointProcessor(BaseDocumentProcessor):
    extensions = ("pptx",)
    file_type = FileType.PPTX
    display_name = "PowerPoint presentation"

    def parse(self, data: bytes, ctx: ProcessingContext) -> Document:
        if not PYTHON_PPTX_AVAILABLE:
            raise CorruptedDocumentError(ctx.filename, "python-pptx is not installed")

        ctx.progress(PipelineStage.PARSING, 0.05, "Opening presentation")
        try:
            presentation = Presentation(io.BytesIO(data))
        except Exception as exc:
            raise CorruptedDocumentError(ctx.filename, str(exc)) from exc

        document = self.new_document(data, ctx)
        slides = list(presentation.slides)
        total = max(1, len(slides))

        for number, slide in enumerate(slides, start=1):
            ctx.progress(
                PipelineStage.EXTRACTING_TEXT,
                0.1 + 0.6 * ((number - 1) / total),
                f"Slide {number} of {total}",
            )
            try:
                page = self._process_slide(slide, number, ctx)
            except Exception as exc:  # noqa: BLE001 - contain per-slide failures
                logger.exception("Failed to process slide %d of %s", number, ctx.filename)
                ctx.warn(f"Slide {number} could not be processed ({type(exc).__name__}).")
                continue
            if page is not None:
                document.pages.append(page)

        document = self.finalize(document, ctx)
        if not document.blocks:
            raise EmptyDocumentError(ctx.filename)
        return document

    # ------------------------------------------------------------------ #
    def _process_slide(self, slide: Any, number: int, ctx: ProcessingContext) -> Optional[Page]:
        page = Page(
            document_id=ctx.document_id,
            session_id=ctx.session_id,
            page_number=number,
            label=f"Slide {number}",
        )

        title = _slide_title(slide)
        if title:
            page.metadata["title"] = title
            block = self.make_text_block(
                ctx,
                page_number=number,
                text=title,
                index=0,
                block_type=BlockType.HEADING,
            )
            if block is not None:
                page.blocks.append(block)

        index = 1
        body_parts: List[str] = []

        for shape in _iter_shapes(slide.shapes):
            index += 1

            # Tables -------------------------------------------------- #
            if getattr(shape, "has_table", False):
                table_block = self._table_block(shape, ctx, number, index, title)
                if table_block is not None:
                    page.blocks.append(table_block)
                continue

            # Native charts ------------------------------------------- #
            if getattr(shape, "has_chart", False):
                chart_block = self._chart_block(shape, ctx, number, index, title)
                if chart_block is not None:
                    page.blocks.append(chart_block)
                continue

            # Pictures ------------------------------------------------ #
            if _is_picture(shape):
                continue  # handled in a second pass, under the image budget

            # Text ---------------------------------------------------- #
            text = _shape_text(shape)
            if text and text != title:
                body_parts.append(text)

        if body_parts:
            block = self.make_text_block(
                ctx,
                page_number=number,
                text="\n".join(body_parts),
                index=1,
                section=title,
            )
            if block is not None:
                page.blocks.append(block)

        self._process_pictures(slide, page, ctx, title or "")
        self._process_notes(slide, page, ctx, title)

        return page if page.blocks else None

    # -- components ------------------------------------------------------ #
    def _table_block(
        self, shape: Any, ctx: ProcessingContext, number: int, index: int, section: Optional[str]
    ) -> Optional[ContentBlock]:
        try:
            rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
        except Exception:
            return None
        data = build_table(rows, has_header=True)
        if data is None:
            return None
        return ContentBlock(
            block_id=self.block_id(ctx, number, index, "table"),
            document_id=ctx.document_id,
            session_id=ctx.session_id,
            page_number=number,
            block_type=BlockType.TABLE,
            source_kind=SourceKind.STRUCTURED,
            text=table_to_text(data),
            table=data,
            parent_section=section,
            order=ctx.next_order(),
        )

    def _chart_block(
        self, shape: Any, ctx: ProcessingContext, number: int, index: int, section: Optional[str]
    ) -> Optional[ContentBlock]:
        """Read a native chart's real numbers from the XML — no model needed."""
        try:
            chart = shape.chart
            chart_type = str(getattr(chart, "chart_type", "")).split(".")[-1].split(" ")[0]
            categories = [str(c) for c in chart.plots[0].categories] if chart.plots else []
            series_lines: List[str] = []
            for series in chart.series:
                values = list(series.values)
                pairs = []
                for position, value in enumerate(values):
                    if value is None:
                        continue
                    label = categories[position] if position < len(categories) else f"#{position + 1}"
                    pairs.append(f"{label}: {_fmt(value)}")
                if pairs:
                    series_lines.append(f"Series '{series.name}': " + ", ".join(pairs))
        except Exception as exc:
            logger.debug("Could not read a native chart on slide %d: %s", number, exc)
            return None

        if not series_lines:
            return None

        title_text = ""
        try:
            if chart.has_title and chart.chart_title.text_frame is not None:
                title_text = clean_text(chart.chart_title.text_frame.text)
        except Exception:
            title_text = ""

        parts = [f"Chart ({chart_type}) on slide {number}."]
        if title_text:
            parts.insert(0, f"Chart title: {title_text}.")
        if categories:
            parts.append("Categories: " + ", ".join(str(c) for c in categories[:30]) + ".")
        parts.extend(series_lines[:20])

        return ContentBlock(
            block_id=self.block_id(ctx, number, index, "chart"),
            document_id=ctx.document_id,
            session_id=ctx.session_id,
            page_number=number,
            block_type=BlockType.CHART,
            # Values come from the file's own data, not from a model reading pixels.
            source_kind=SourceKind.STRUCTURED,
            text="\n".join(parts),
            parent_section=section,
            order=ctx.next_order(),
            metadata={"chart_type": chart_type, "visual_title": title_text},
        )

    def _process_pictures(
        self, slide: Any, page: Page, ctx: ProcessingContext, context_text: str
    ) -> None:
        if not ctx.settings.vision.enabled:
            return
        analysed = 0
        for index, shape in enumerate(_iter_shapes(slide.shapes)):
            if analysed >= MAX_IMAGES_PER_SLIDE:
                break
            if not _is_picture(shape):
                continue
            try:
                image_bytes = shape.image.blob
            except Exception:
                continue
            if not image_bytes or len(image_bytes) < MIN_IMAGE_BYTES:
                continue
            if is_probably_decorative(
                image_bytes, min_pixels=ctx.settings.vision.min_image_pixels
            ):
                continue

            block = self.process_visual(
                ctx,
                image_bytes,
                page_number=page.page_number,
                index=300 + index,
                context_text=context_text,
                origin="embedded",
            )
            if block is not None:
                page.blocks.append(block)
                analysed += 1

    def _process_notes(
        self, slide: Any, page: Page, ctx: ProcessingContext, section: Optional[str]
    ) -> None:
        try:
            if not slide.has_notes_slide:
                return
            notes = clean_text(slide.notes_slide.notes_text_frame.text)
        except Exception:
            return
        if not notes:
            return
        block = self.make_text_block(
            ctx,
            page_number=page.page_number,
            text=notes,
            index=900,
            block_type=BlockType.SPEAKER_NOTES,
            section=section,
        )
        if block is not None:
            page.blocks.append(block)


# --------------------------------------------------------------------------- #
def _iter_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    """Walk shapes recursively so grouped diagram parts are not missed."""
    for shape in shapes:
        if MSO_SHAPE_TYPE is not None and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                yield from _iter_shapes(shape.shapes)
                continue
            except Exception:
                pass
        yield shape


def _is_picture(shape: Any) -> bool:
    if MSO_SHAPE_TYPE is None:
        return False
    try:
        return shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE)
    except Exception:
        return False


def _shape_text(shape: Any) -> str:
    try:
        if not getattr(shape, "has_text_frame", False):
            return ""
        return clean_text(shape.text_frame.text)
    except Exception:
        return ""


def _slide_title(slide: Any) -> str:
    try:
        placeholder = slide.shapes.title
        if placeholder is not None:
            return clean_text(placeholder.text)
    except Exception:
        pass
    return ""


def _fmt(value: Any) -> str:
    try:
        number = float(value)
    except (TypeError, ValueError):
        return str(value)
    if number == int(number) and abs(number) < 1e15:
        return str(int(number))
    return f"{number:g}"
